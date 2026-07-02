from __future__ import annotations

import math
import os
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(
    os.environ.get(
        "CO2_PROJECT_ROOT",
        r"E:\BaiduSyncdisk\desktop_5_15\01_科研项目\GAM_CO2_SctO2_4_19_2026",
    )
)
SOURCE_CSV = Path(
    os.environ.get(
        "SUBGROUP_ABS_SOURCE_CSV",
        PROJECT_ROOT
        / "hpc_r_format_outputs"
        / "subgroup_heatmap_table1_table2_20260701"
        / "subgroup_boot200_endpoint_effects_raw.csv",
    )
)
OUT_DIR = Path(
    os.environ.get(
        "SUBGROUP_ABS_OUT_DIR",
        PROJECT_ROOT
        / "hpc_r_format_outputs"
        / "subgroup_heatmap_table1_table2_20260701"
        / "native_abs_editable_20260702",
    )
)

ROW_ORDER = [
    "Overall",
    "Age <70 year",
    "Age ≥70 year",
    "Female",
    "Male",
    "Preop BP\n<140/90 mmHg",
    "Preop BP\n≥140/90 mmHg",
]
Y_ORDER = ["Left SctO₂", "Right SctO₂", "SftO₂"]
X_ORDER = ["ET_CO2", "FiO2_new", "TEMP"]
X_LABEL = {"ET_CO2": "EtCO₂", "FiO2_new": "FiO₂", "TEMP": "TEMP"}
Y_LABEL = {"Left SctO₂": "Left SctO₂", "Right SctO₂": "Right SctO₂", "SftO₂": "SftO₂"}

FONT = "Aptos"
FILL_MIN = 0.0
FILL_MAX = 7.0
LOW_COLOR = (247, 247, 247)
HIGH_COLOR = (253, 80, 2)


def interp_color(v: float) -> RGBColor:
    t = (max(FILL_MIN, min(FILL_MAX, v)) - FILL_MIN) / (FILL_MAX - FILL_MIN)
    rgb = tuple(round(LOW_COLOR[i] + t * (HIGH_COLOR[i] - LOW_COLOR[i])) for i in range(3))
    return RGBColor(*rgb)


def add_text(slide, text, left, top, width, height, size=10, bold=False, color=(0, 0, 0), rotate=0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.rotation = rotate
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.word_wrap = False
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_rect(slide, left, top, width, height, fill, line=(255, 255, 255), line_width=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_width == 0:
        shp.line.fill.background()
        shp.line.width = Pt(0)
    else:
        shp.line.color.rgb = RGBColor(*line)
        shp.line.width = Pt(line_width)
    try:
        shp.shadow.inherit = False
        shp.shadow.visible = False
    except Exception:
        pass
    return shp


def add_line(slide, x1, y1, x2, y2, color=(95, 95, 95), width=0.7):
    line = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(width)
    return line


def add_border(slide, left, top, width, height, color=(95, 95, 95), line_width=0.8):
    add_line(slide, left, top, left + width, top, color=color, width=line_width)
    add_line(slide, left, top + height, left + width, top + height, color=color, width=line_width)
    add_line(slide, left, top, left, top + height, color=color, width=line_width)
    add_line(slide, left + width, top, left + width, top + height, color=color, width=line_width)


def load_abs_summary() -> pd.DataFrame:
    raw = pd.read_csv(SOURCE_CSV)
    raw["abs_full_5_95_change"] = raw["full_5_95_change"].abs()
    out = (
        raw.groupby(["subgroup", "ycol", "xvar", "window_lo", "window_hi"], as_index=False)
        .agg(
            mean=("abs_full_5_95_change", "mean"),
            lo95=("abs_full_5_95_change", lambda s: s.quantile(0.025)),
            hi95=("abs_full_5_95_change", lambda s: s.quantile(0.975)),
        )
    )
    out["value_label"] = out["mean"].map(lambda x: "0.00" if abs(x) < 0.005 else f"{x:.2f}")
    return out


def build_ppt(summary: pd.DataFrame, ppt_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    cell_w = 0.72
    cell_h = 0.58
    block_gap = 0.34
    grid_top = 0.82
    grid_left = 2.70
    block_w = cell_w * len(X_ORDER)
    block_h = cell_h * len(ROW_ORDER)

    label_left = 0.55
    label_w = 1.95
    title_h = 0.30

    for i, row_label in enumerate(ROW_ORDER):
        y = grid_top + i * cell_h
        add_text(slide, row_label, label_left, y + 0.02, label_w, cell_h - 0.04, size=12.0)

    for b, ycol in enumerate(Y_ORDER):
        x0 = grid_left + b * (block_w + block_gap)
        add_text(slide, Y_LABEL[ycol], x0, grid_top - 0.45, block_w, title_h, size=15.0)

        block = summary[summary["ycol"].eq(ycol)]
        for i, row_label in enumerate(ROW_ORDER):
            for j, xvar in enumerate(X_ORDER):
                cell = block[block["subgroup"].eq(row_label) & block["xvar"].eq(xvar)]
                if cell.empty:
                    value = math.nan
                    label = ""
                else:
                    value = float(cell.iloc[0]["mean"])
                    label = str(cell.iloc[0]["value_label"])
                x = x0 + j * cell_w
                y = grid_top + i * cell_h
                fill = interp_color(0.0 if math.isnan(value) else value)
                add_rect(slide, x, y, cell_w, cell_h, fill, line=(255, 255, 255), line_width=1.0)
                text_color = (255, 255, 255) if value >= 5.3 else (0, 0, 0)
                add_text(slide, label, x, y, cell_w, cell_h, size=9.6, color=text_color)

        add_border(slide, x0, grid_top, block_w, block_h)
        for j, xvar in enumerate(X_ORDER):
            cx = x0 + j * cell_w + cell_w / 2
            add_line(slide, cx, grid_top + block_h, cx, grid_top + block_h + 0.08, width=0.7)
            add_text(slide, X_LABEL[xvar], x0 + j * cell_w, grid_top + block_h + 0.14, cell_w, 0.24, size=12.0)

    bar_x = grid_left + 3 * block_w + 2 * block_gap + 0.22
    bar_y = grid_top
    bar_w = 0.14
    bar_h = block_h
    steps = 80
    for k in range(steps):
        frac = k / steps
        val_low = FILL_MIN + frac * (FILL_MAX - FILL_MIN)
        y = bar_y + bar_h * (1 - (k + 1) / steps)
        fill = interp_color(val_low)
        add_rect(slide, bar_x, y, bar_w, bar_h / steps + 0.002, fill, line=(255, 255, 255), line_width=0)
    add_border(slide, bar_x, bar_y, bar_w, bar_h, color=(51, 51, 51), line_width=0.8)

    for tick in range(0, 8):
        ty = bar_y + bar_h * (1 - tick / FILL_MAX)
        add_line(slide, bar_x + bar_w, ty, bar_x + bar_w + 0.07, ty, color=(51, 51, 51), width=0.7)
        add_text(slide, f"{tick:.1f}", bar_x + bar_w + 0.10, ty - 0.08, 0.35, 0.16, size=9.5)

    add_text(
        slide,
        "Tissue Oxygenation Change (%)",
        bar_x + 0.12,
        bar_y + bar_h / 2 - 0.15,
        bar_h,
        0.30,
        size=11.5,
        rotate=90,
    )

    ppt_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(ppt_path)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    summary = load_abs_summary()
    summary.to_csv(OUT_DIR / "subgroup_table2_full_5_95_abs_change_mean_95ci.csv", index=False)
    build_ppt(summary, OUT_DIR / "Figure4_heatmap_abs_native_editable.pptx")

    short_path = Path(r"E:\BaiduSyncdisk\desktop_5_15\Figure4_heatmap_abs_native_editable.pptx")
    build_ppt(summary, short_path)
    print(OUT_DIR / "Figure4_heatmap_abs_native_editable.pptx")
    print(short_path)


if __name__ == "__main__":
    main()
