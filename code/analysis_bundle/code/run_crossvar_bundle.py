#!/usr/bin/env python3
import argparse
import logging
from pathlib import Path
from datetime import datetime

import numpy as np
import pandas as pd
import yaml

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import TwoSlopeNorm

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None
    Inches = None
    Pt = None


def resolve_path(base, p):
    p = Path(p)
    if p.is_absolute():
        return p
    return (base / p).resolve()


def setup_logger(log_fp):
    log_fp.parent.mkdir(parents=True, exist_ok=True)
    logger = logging.getLogger("crossvar_bundle")
    logger.setLevel(logging.INFO)
    logger.handlers = []
    fmt = logging.Formatter("%(asctime)s | %(levelname)s | %(message)s")
    ch = logging.StreamHandler()
    ch.setFormatter(fmt)
    logger.addHandler(ch)
    fh = logging.FileHandler(log_fp)
    fh.setFormatter(fmt)
    logger.addHandler(fh)
    return logger


def safe_interp(x, y, xq):
    return float(np.interp(float(xq), x, y))


def clean_sign(arr, eps):
    s = np.sign(arr)
    s[np.abs(arr) <= eps] = 0
    for i in range(1, len(s)):
        if s[i] == 0:
            s[i] = s[i - 1]
    for i in range(len(s) - 2, -1, -1):
        if s[i] == 0:
            s[i] = s[i + 1]
    return s


def find_turnpoint(x, slope):
    if len(x) < 3:
        return np.nan
    eps = np.nanquantile(np.abs(slope), 0.1)
    s = clean_sign(slope.copy(), eps)
    idx = np.where(s[:-1] * s[1:] < 0)[0]
    if len(idx) == 0:
        return np.nan
    i = int(idx[0])
    return float((x[i] + x[i + 1]) / 2.0)


def find_plateau(x, slope):
    if len(x) < 10:
        return (np.nan, np.nan)
    th = float(np.nanquantile(np.abs(slope), 0.25))
    mask = np.abs(slope) <= th
    best = (0, -1)
    start = None
    for i, v in enumerate(mask):
        if v and start is None:
            start = i
        if (not v or i == len(mask) - 1) and start is not None:
            end = i if v and i == len(mask) - 1 else i - 1
            if end - start > best[1] - best[0]:
                best = (start, end)
            start = None
    if best[1] < best[0]:
        return (np.nan, np.nan)
    return (float(x[best[0]]), float(x[best[1]]))


def load_curve(curve_fp):
    df = pd.read_csv(curve_fp)
    need = ["x", "pred_mean", "pred_lo_2.5", "pred_hi_97.5"]
    for c in need:
        if c not in df.columns:
            raise ValueError("missing column {} in {}".format(c, curve_fp))
    x = pd.to_numeric(df["x"], errors="coerce").to_numpy(dtype=float)
    y = pd.to_numeric(df["pred_mean"], errors="coerce").to_numpy(dtype=float)
    lo = pd.to_numeric(df["pred_lo_2.5"], errors="coerce").to_numpy(dtype=float)
    hi = pd.to_numeric(df["pred_hi_97.5"], errors="coerce").to_numpy(dtype=float)
    good = np.isfinite(x) & np.isfinite(y) & np.isfinite(lo) & np.isfinite(hi)
    x, y, lo, hi = x[good], y[good], lo[good], hi[good]
    order = np.argsort(x)
    x, y, lo, hi = x[order], y[order], lo[order], hi[order]

    xvar = str(df["xvar"].iloc[0]) if "xvar" in df.columns else curve_fp.parent.name
    ycol = str(df["ycol"].iloc[0]) if "ycol" in df.columns else curve_fp.parents[3].name
    sec = str(df["sec"].iloc[0]) if "sec" in df.columns else "1"
    n_sample = str(df["n_sample"].iloc[0]) if "n_sample" in df.columns else ""

    return {
        "x": x,
        "y": y,
        "lo": lo,
        "hi": hi,
        "xvar": xvar,
        "ycol": ycol,
        "sec": sec,
        "n_sample": n_sample,
    }


def load_boot_matrix(curve_fp, x):
    matrix_fp = Path(str(curve_fp).replace("_curve_boot.csv", "_boot_raw_curve_matrix.csv"))
    if not matrix_fp.exists():
        return None, None
    mat_df = pd.read_csv(matrix_fp)
    mat = mat_df.to_numpy(dtype=float)
    if mat.ndim != 2 or mat.shape[0] < 5:
        return None, matrix_fp
    if mat.shape[1] != len(x):
        x_mat = np.linspace(float(np.nanmin(x)), float(np.nanmax(x)), mat.shape[1])
    else:
        x_mat = x
    return (mat, x_mat), matrix_fp


def delta_ci_from_boot(mat, x_mat, x0, x1):
    d = []
    for i in range(mat.shape[0]):
        row = mat[i, :]
        if not np.all(np.isfinite(row)):
            continue
        y0 = safe_interp(x_mat, row, x0)
        y1 = safe_interp(x_mat, row, x1)
        d.append(y1 - y0)
    if len(d) < 10:
        return np.nan, np.nan, np.nan
    d = np.array(d, dtype=float)
    return float(np.quantile(d, 0.025)), float(np.quantile(d, 0.975)), float(np.mean(d))


def delta_ci_approx_from_pointwise(x, lo, hi, x0, x1, delta):
    y0_lo, y0_hi = safe_interp(x, lo, x0), safe_interp(x, hi, x0)
    y1_lo, y1_hi = safe_interp(x, lo, x1), safe_interp(x, hi, x1)
    se0 = max((y0_hi - y0_lo) / (2.0 * 1.96), 1e-9)
    se1 = max((y1_hi - y1_lo) / (2.0 * 1.96), 1e-9)
    se = np.sqrt(se0 ** 2 + se1 ** 2)
    return float(delta - 1.96 * se), float(delta + 1.96 * se)


def clinical_window_effects(x, y, step, window_lo, window_hi, n_segments):
    if not np.isfinite(window_lo) or not np.isfinite(window_hi) or window_hi <= window_lo:
        raise ValueError("invalid clinical effect window")
    edges = np.linspace(float(window_lo), float(window_hi), int(n_segments) + 1)
    rows = []
    effects = []
    for i in range(int(n_segments)):
        x0 = float(edges[i])
        x1 = float(edges[i + 1])
        if x1 <= x0:
            continue
        y0 = safe_interp(x, y, x0)
        y1 = safe_interp(x, y, x1)
        slope = (y1 - y0) / (x1 - x0)
        effect = slope * float(step)
        rows.append({
            "bin_idx": i,
            "x0": x0,
            "x1": x1,
            "x_mid": float((x0 + x1) / 2.0),
            "slope_per_unit": float(slope),
            "clinical_step_effect": float(effect),
        })
        effects.append(float(effect))
    if not effects:
        return np.nan, np.nan, np.nan, rows
    effects = np.asarray(effects, dtype=float)
    return (
        float(np.nanmedian(effects)),
        float(np.nanquantile(effects, 0.25)),
        float(np.nanquantile(effects, 0.75)),
        rows,
    )


def compute_one(curve_fp, step, q_window, window_row=None, n_segments=20):
    c = load_curve(curve_fp)
    x, y, lo, hi = c["x"], c["y"], c["lo"], c["hi"]
    xmin, xmax = float(np.nanmin(x)), float(np.nanmax(x))

    if window_row is not None:
        window_lo = float(window_row["q05"])
        window_hi = float(window_row["q95"])
    else:
        window_lo = float(np.quantile(x, q_window[0]))
        window_hi = float(np.quantile(x, q_window[1]))
    window_lo = max(xmin, window_lo)
    window_hi = min(xmax, window_hi)

    delta, iqr_lo, iqr_hi, segment_rows = clinical_window_effects(
        x=x, y=y, step=step, window_lo=window_lo, window_hi=window_hi, n_segments=n_segments
    )
    ci_lo, ci_hi = iqr_lo, iqr_hi
    delta_boot_mean = np.nan
    ci_source = "pred_mean_5_95_20_segments_iqr"
    x0, x1 = window_lo, window_hi

    slope = np.gradient(y, x)
    slope_at_median = safe_interp(x, slope, float((window_lo + window_hi) / 2.0))
    mask = (x >= window_lo) & (x <= window_hi)
    mean_abs_slope_q10_q90 = float(np.nanmean(np.abs(slope[mask]))) if np.any(mask) else np.nan

    turnpoint_x = find_turnpoint(x, slope)
    plateau_x0, plateau_x1 = find_plateau(x, slope)

    slope_bins = []
    for rr in segment_rows:
        slope_bins.append(float(rr["clinical_step_effect"]))

    return {
        "xvar": c["xvar"],
        "ycol": c["ycol"],
        "sec": c["sec"],
        "n_sample": c["n_sample"],
        "curve_fp": str(curve_fp),
        "boot_matrix_fp": "",
        "x_min": xmin,
        "x_max": xmax,
        "x0": x0,
        "x1": x1,
        "effect_window_lo": window_lo,
        "effect_window_hi": window_hi,
        "effect_window_source": "model_input_q05_q95" if window_row is not None else "curve_x_quantiles",
        "effect_n_segments": int(n_segments),
        "clinical_step": float(step),
        "delta_rso2_clinical_step": delta,
        "delta_rso2_ci_lo": ci_lo,
        "delta_rso2_ci_hi": ci_hi,
        "delta_rso2_iqr_lo": iqr_lo,
        "delta_rso2_iqr_hi": iqr_hi,
        "delta_rso2_summary": "median_iqr_across_20_segments",
        "delta_rso2_boot_mean": delta_boot_mean,
        "ci_source": ci_source,
        "slope_at_median": slope_at_median,
        "mean_abs_slope_q10_q90": mean_abs_slope_q10_q90,
        "turnpoint_x": turnpoint_x,
        "plateau_x0": plateau_x0,
        "plateau_x1": plateau_x1,
        "slope_bins": slope_bins,
        "segment_rows": segment_rows,
        "curve_x": x,
        "curve_y": y,
        "curve_lo": lo,
        "curve_hi": hi,
    }


def plot_bar(df, cfg, out_png, out_pdf):
    channels = cfg["channels"]
    variables = cfg["variables"]
    colors = cfg["colors"]

    width = max(16, 1.4 * len(variables) * max(1, len(channels)))
    fig, axes = plt.subplots(1, len(channels), figsize=(width, 5), sharey=True)
    if len(channels) == 1:
        axes = [axes]

    for ax, ch in zip(axes, channels):
        d = df[(df["ycol"] == ch) & (df["status"] == "ok")].copy()
        x_pos = np.arange(len(variables))
        y = []
        yerr_low = []
        yerr_high = []
        for v in variables:
            r = d[d["xvar"] == v]
            if r.empty:
                y.append(np.nan)
                yerr_low.append(np.nan)
                yerr_high.append(np.nan)
            else:
                rr = r.iloc[0]
                y.append(float(rr["delta_rso2_clinical_step"]))
                yerr_low.append(float(rr["delta_rso2_clinical_step"]) - float(rr["delta_rso2_ci_lo"]))
                yerr_high.append(float(rr["delta_rso2_ci_hi"]) - float(rr["delta_rso2_clinical_step"]))

        for i, v in enumerate(variables):
            if np.isnan(y[i]):
                ax.bar(i, 0.0, color="#cccccc", edgecolor="black", hatch="//")
                ax.text(i, 0.1, "missing", ha="center", va="bottom", fontsize=8, rotation=90)
            else:
                ax.bar(i, y[i], color=colors.get(v, "#888888"), edgecolor="black")
                ax.errorbar(i, y[i], yerr=[[yerr_low[i]], [yerr_high[i]]], color="black", capsize=3, lw=1)

        ax.axhline(0, color="black", lw=0.8)
        ax.set_xticks(x_pos)
        ax.set_xticklabels(variables, rotation=35, ha="right")
        ax.set_title(ch)
        ax.grid(alpha=0.2, axis="y")

    fig.suptitle("Figure A: Median clinical-step effect across the model-input 5%-95% window", fontsize=13)
    axes[0].set_ylabel("Delta rSO2, median (IQR)")
    ylims = cfg["plot"]["y_limits"]
    yticks = cfg["plot"]["y_ticks"]
    if ylims and len(ylims) == 2:
        axes[0].set_ylim(float(ylims[0]) - 70.0, float(ylims[1]) - 70.0)
    for ax in axes:
        ax.set_yticks(np.arange(-4, 5, 1))

    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(out_png, dpi=int(cfg["plot"]["dpi"]))
    fig.savefig(out_pdf)
    plt.close(fig)


def plot_heatmap(df, cfg, out_png, out_pdf):
    channels = cfg["channels"]
    variables = cfg["variables"]
    n_segments = int(cfg["plot"].get("n_segments", 20))
    bin_labels = ["S{:02d}".format(i + 1) for i in range(n_segments)]

    mats = {}
    all_vals = []
    for v in variables:
        mat = np.full((len(channels), n_segments), np.nan)
        for i, ch in enumerate(channels):
            r = df[(df["status"] == "ok") & (df["xvar"] == v) & (df["ycol"] == ch)]
            if r.empty:
                continue
            bins = r.iloc[0]["slope_bins"]
            vals = [float(x) if np.isfinite(x) else np.nan for x in bins]
            mat[i, :] = vals
            all_vals.extend([x for x in vals if np.isfinite(x)])
        mats[v] = mat

    vmax = float(np.nanmax(np.abs(all_vals))) if len(all_vals) else 1.0
    vmax = max(vmax, 1e-6)
    norm = TwoSlopeNorm(vmin=-vmax, vcenter=0.0, vmax=vmax)

    width = max(18, 1.8 * len(variables))
    fig, axes = plt.subplots(1, len(variables), figsize=(width, 5), sharey=True)
    if len(variables) == 1:
        axes = [axes]

    for ax, v in zip(axes, variables):
        im = ax.imshow(mats[v], aspect="auto", cmap="coolwarm", norm=norm)
        ax.set_title(v)
        ax.set_xticks(np.arange(n_segments))
        ax.set_xticklabels(bin_labels, rotation=45, ha="right", fontsize=8)
        ax.set_yticks(np.arange(len(channels)))
        ax.set_yticklabels(channels)

    cbar = fig.colorbar(im, ax=axes, fraction=0.02, pad=0.02)
    cbar.set_label("Clinical-step effect per segment")
    fig.suptitle("Figure B: 20-segment clinical-step effects across the model-input 5%-95% window", fontsize=13)
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(out_png, dpi=int(cfg["plot"]["dpi"]))
    fig.savefig(out_pdf)
    plt.close(fig)


def plot_threshold(df, cfg, out_png, out_pdf):
    channels = cfg["channels"]
    variables = cfg["variables"]

    width = max(18, 1.8 * len(variables))
    height = max(12, 3.5 * len(channels))
    fig, axes = plt.subplots(len(channels), len(variables), figsize=(width, height), sharey=True)
    if len(channels) == 1 and len(variables) == 1:
        axes = np.array([[axes]])
    elif len(channels) == 1:
        axes = np.array([axes])
    elif len(variables) == 1:
        axes = np.array([[a] for a in axes])

    for i, ch in enumerate(channels):
        for j, v in enumerate(variables):
            ax = axes[i, j]
            r = df[(df["status"] == "ok") & (df["xvar"] == v) & (df["ycol"] == ch)]
            if r.empty:
                ax.text(0.5, 0.5, "missing", transform=ax.transAxes, ha="center", va="center")
                ax.set_title("{} | {}".format(ch, v))
                continue

            rr = r.iloc[0]
            x = rr["curve_x"]
            y = rr["curve_y"]
            lo = rr["curve_lo"]
            hi = rr["curve_hi"]

            ax.plot(x, y, color=cfg["colors"].get(v, "#444444"), lw=2)
            ax.fill_between(x, lo, hi, color=cfg["colors"].get(v, "#444444"), alpha=0.2)

            if np.isfinite(rr["turnpoint_x"]):
                ax.axvline(float(rr["turnpoint_x"]), color="black", linestyle="--", lw=1)

            if np.isfinite(rr["plateau_x0"]) and np.isfinite(rr["plateau_x1"]):
                ax.axvspan(float(rr["plateau_x0"]), float(rr["plateau_x1"]), color="gray", alpha=0.15)

            ax.set_title("{} | {}".format(ch, v), fontsize=10)
            ax.grid(alpha=0.2)

    for j, v in enumerate(variables):
        axes[-1, j].set_xlabel(v)
    for i, ch in enumerate(channels):
        axes[i, 0].set_ylabel(ch)

    fig.suptitle("Figure C: Threshold / turning-point view with CI bands", fontsize=13)
    fig.tight_layout(rect=[0, 0, 1, 0.96])
    fig.savefig(out_png, dpi=int(cfg["plot"]["dpi"]))
    fig.savefig(out_pdf)
    plt.close(fig)


def add_title(slide, text, left=0.5, top=0.2, width=12.0, height=0.6, size=28):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True


def add_body(slide, text, left=0.6, top=1.0, width=12.0, height=1.5, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)


def build_ppt(cfg, run_meta, fig_paths, table_df, out_ppt):
    if Presentation is None:
        return False
    prs = Presentation()

    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s0, cfg["ppt"]["title"])
    add_body(s0, cfg["ppt"]["subtitle"], top=1.1, size=14)
    add_body(
        s0,
        "Run: {} | Model: {} | Scope: {}\\nSource: {}".format(
            run_meta["run_id"], run_meta["model_label"], run_meta["scope"], run_meta["result_dir"]
        ),
        top=2.0,
        size=12,
    )

    slides = [
        ("Main effect comparison", fig_paths["figure_a_png"]),
        ("Slope heatmap", fig_paths["figure_b_png"]),
        ("Threshold and turning-point view", fig_paths["figure_c_png"]),
    ]

    for title, img_fp in slides:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, title, size=24)
        s.shapes.add_picture(str(img_fp), Inches(0.5), Inches(0.9), width=Inches(12.3))

    # Appendix summary
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "Appendix: key numeric summary", size=22)
    lines = []
    cols = ["ycol", "xvar", "delta_rso2_clinical_step", "delta_rso2_iqr_lo", "delta_rso2_iqr_hi", "effect_window_lo", "effect_window_hi"]
    ok = table_df[table_df["status"] == "ok"].copy()
    for _, r in ok[cols].iterrows():
        lines.append(
            "{} | {}: median Delta={:.3f} (IQR {:.3f},{:.3f}), window={:.3f}-{:.3f}".format(
                r["ycol"], r["xvar"], r["delta_rso2_clinical_step"], r["delta_rso2_iqr_lo"], r["delta_rso2_iqr_hi"], r["effect_window_lo"], r["effect_window_hi"]
            )
        )
    add_body(s4, "\\n".join(lines[:12]), top=0.9, height=6.0, size=11)

    out_ppt.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_ppt))
    return True


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", default="code/config.yaml")
    parser.add_argument("--run-id", default="")
    args = parser.parse_args()

    workspace = Path(__file__).resolve().parents[1]
    cfg_fp = resolve_path(workspace, args.config)
    with cfg_fp.open("r", encoding="utf-8") as f:
        cfg = yaml.safe_load(f)

    out_logs = resolve_path(workspace, cfg["output"]["logs_dir"])
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    logger = setup_logger(out_logs / ("run_crossvar_bundle_{}.log".format(stamp)))

    for k in ["figures_dir", "tables_dir", "ppt_dir"]:
        resolve_path(workspace, cfg["output"][k]).mkdir(parents=True, exist_ok=True)

    run_catalog_fp = resolve_path(workspace, cfg["run_catalog"])
    run_catalog = pd.read_csv(run_catalog_fp)

    run_id = args.run_id.strip() or str(cfg["primary_run_id"])
    output_tag = str(cfg.get("output_tag", run_id)).strip() or run_id
    run_row = run_catalog.loc[run_catalog["run_id"] == run_id]
    if run_row.empty:
        raise SystemExit("run_id not found: {}".format(run_id))
    run_row = run_row.iloc[0].to_dict()

    result_dir = Path(run_row["result_dir"])
    if not result_dir.exists():
        repo_root = workspace.parents[1]
        local_result_dir = repo_root / "results" / "model_runs" / result_dir.name
        if local_result_dir.exists():
            result_dir = local_result_dir
        else:
            raise SystemExit("result_dir missing: {}".format(result_dir))

    logger.info("run_id=%s", run_id)
    logger.info("result_dir=%s", result_dir)

    variables = list(cfg["variables"])
    channels = list(cfg["channels"])
    steps = dict(cfg["clinical_steps"])
    q_window = tuple(cfg["plot"]["q_window"])
    n_segments = int(cfg["plot"].get("n_segments", 20))

    window_lookup = {}
    quant_fp_raw = str(cfg.get("model_input_quantiles", "")).strip()
    if quant_fp_raw:
        quant_fp = resolve_path(workspace, quant_fp_raw)
        if quant_fp.exists():
            quant_df = pd.read_csv(quant_fp)
            for _, rr in quant_df.iterrows():
                window_lookup[(str(rr["ycol"]), str(rr["xvar"]))] = rr
            logger.info("model_input_quantiles=%s rows=%s", quant_fp, len(quant_df))
        else:
            logger.warning("model_input_quantiles missing: %s", quant_fp)

    curve_files = sorted(result_dir.rglob("*_curve_boot.csv"))
    curve_files = [fp for fp in curve_files if fp.name.endswith("_{}_curve_boot.csv".format(cfg["plot_mode"]))]

    key_to_fp = {}
    for fp in curve_files:
        try:
            info = load_curve(fp)
        except Exception:
            continue
        key = (info["ycol"], info["xvar"])
        if info["ycol"] in channels and info["xvar"] in variables:
            key_to_fp[key] = fp

    rows = []
    slope_rows = []

    for ch in channels:
        for v in variables:
            key = (ch, v)
            fp = key_to_fp.get(key)
            if fp is None:
                wrow = window_lookup.get((ch, v))
                rows.append({
                    "run_id": run_id,
                    "model_label": run_row["model_label"],
                    "scope": run_row["scope"],
                    "ycol": ch,
                    "xvar": v,
                    "clinical_step": float(steps[v]),
                    "effect_window_lo": float(wrow["q05"]) if wrow is not None else np.nan,
                    "effect_window_hi": float(wrow["q95"]) if wrow is not None else np.nan,
                    "effect_window_source": "model_input_q05_q95" if wrow is not None else "",
                    "effect_n_segments": int(n_segments),
                    "status": "missing",
                    "curve_fp": "",
                })
                continue

            rec = compute_one(
                fp,
                step=float(steps[v]),
                q_window=q_window,
                window_row=window_lookup.get((ch, v)),
                n_segments=n_segments,
            )
            rec.update({
                "run_id": run_id,
                "model_label": run_row["model_label"],
                "scope": run_row["scope"],
                "status": "ok",
            })
            rows.append(rec)

            for seg in rec["segment_rows"]:
                slope_rows.append({
                    "run_id": run_id,
                    "ycol": ch,
                    "xvar": v,
                    "bin_idx": int(seg["bin_idx"]),
                    "bin_label": "S{:02d}".format(int(seg["bin_idx"]) + 1),
                    "x0": seg["x0"],
                    "x1": seg["x1"],
                    "x_mid": seg["x_mid"],
                    "slope_per_unit": seg["slope_per_unit"],
                    "clinical_step_effect": seg["clinical_step_effect"],
                })

    df = pd.DataFrame(rows)
    df_slope = pd.DataFrame(slope_rows)

    out_tables = resolve_path(workspace, cfg["output"]["tables_dir"])
    # Save detailed table without arrays first
    drop_cols = ["slope_bins", "segment_rows", "curve_x", "curve_y", "curve_lo", "curve_hi"]
    cols_keep = [c for c in df.columns if c not in drop_cols]
    summary_fp = out_tables / ("crossvar_effect_summary_{}.csv".format(output_tag))
    df[cols_keep].to_csv(summary_fp, index=False)

    by_channel = (
        df[df["status"] == "ok"]
        .pivot_table(index=["run_id", "ycol"], columns="xvar", values="delta_rso2_clinical_step", aggfunc="first")
        .reset_index()
    )
    by_channel_fp = out_tables / ("crossvar_effect_by_channel_{}.csv".format(output_tag))
    by_channel.to_csv(by_channel_fp, index=False)

    slope_fp = out_tables / ("crossvar_slope_bins_{}.csv".format(output_tag))
    df_slope.to_csv(slope_fp, index=False)

    out_fig = resolve_path(workspace, cfg["output"]["figures_dir"])
    fig_a_png, fig_a_pdf = out_fig / ("figure_A_delta_bar_{}.png".format(output_tag)), out_fig / ("figure_A_delta_bar_{}.pdf".format(output_tag))
    fig_b_png, fig_b_pdf = out_fig / ("figure_B_slope_heatmap_{}.png".format(output_tag)), out_fig / ("figure_B_slope_heatmap_{}.pdf".format(output_tag))
    fig_c_png, fig_c_pdf = out_fig / ("figure_C_threshold_turning_{}.png".format(output_tag)), out_fig / ("figure_C_threshold_turning_{}.pdf".format(output_tag))

    plot_bar(df, cfg, fig_a_png, fig_a_pdf)
    plot_heatmap(df, cfg, fig_b_png, fig_b_pdf)
    plot_threshold(df, cfg, fig_c_png, fig_c_pdf)

    fig_index = pd.DataFrame([
        {"figure_id": "A", "title": "Clinical-step comparable effect", "png": str(fig_a_png), "pdf": str(fig_a_pdf)},
        {"figure_id": "B", "title": "Slope heatmap", "png": str(fig_b_png), "pdf": str(fig_b_pdf)},
        {"figure_id": "C", "title": "Threshold/turning-point view", "png": str(fig_c_png), "pdf": str(fig_c_pdf)},
    ])
    fig_index_fp = out_tables / ("figure_index_{}.csv".format(output_tag))
    fig_index.to_csv(fig_index_fp, index=False)

    out_ppt_name = str(cfg["ppt"]["filename"])
    if out_ppt_name.lower().endswith(".pptx"):
        out_ppt_name = out_ppt_name[:-5] + "_" + output_tag + ".pptx"
    else:
        out_ppt_name = out_ppt_name + "_" + output_tag + ".pptx"
    out_ppt = resolve_path(workspace, cfg["output"]["ppt_dir"]) / out_ppt_name
    fig_paths = {
        "figure_a_png": fig_a_png,
        "figure_b_png": fig_b_png,
        "figure_c_png": fig_c_png,
    }
    ppt_written = build_ppt(cfg, run_row, fig_paths, df[cols_keep], out_ppt)

    checks = {
        "n_expected": len(channels) * len(variables),
        "n_rows": int(len(df)),
        "n_missing": int((df["status"] == "missing").sum()),
        "summary_csv": str(summary_fp),
        "by_channel_csv": str(by_channel_fp),
        "ppt": str(out_ppt) if ppt_written else "",
    }
    pd.DataFrame([checks]).to_csv(out_tables / ("run_checks_{}.csv".format(output_tag)), index=False)

    logger.info("done | rows=%s missing=%s", checks["n_rows"], checks["n_missing"])
    logger.info("summary=%s", summary_fp)
    logger.info("ppt=%s", out_ppt if ppt_written else "skipped: python-pptx not installed")


if __name__ == "__main__":
    main()
