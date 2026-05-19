#!/usr/bin/env python3
"""
v1.1 - 样本量敏感性汇总与 PPTX 报告生成脚本
改动：
  - 支持 python-pptx 自动将效应量折线图与曲线重叠对比图生成专业级幻灯片。
  - 增加 PPT 标题和图表排版辅助函数。
基于：原始 summarize_modelB_n_sweep_etco2.py
解决问题：增加自动 PPT 报告生成，提升结果展现的直观性。
"""
import argparse
import re
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

CHANNELS = ["rSO2_Ch1", "rSO2_Ch2", "rSO2_Ch3"]

def add_title(slide, text, size=24):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.name = "Arial"

def add_body(slide, text, top=1.5, size=14, height=4.5):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.33), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Arial"



def interp(x, y, xq):
    return float(np.interp(float(xq), x, y))


def load_curve(curve_fp: Path):
    df = pd.read_csv(curve_fp)
    x = pd.to_numeric(df["x"], errors="coerce").to_numpy(dtype=float)
    y = pd.to_numeric(df["pred_mean"], errors="coerce").to_numpy(dtype=float)
    lo = pd.to_numeric(df["pred_lo_2.5"], errors="coerce").to_numpy(dtype=float)
    hi = pd.to_numeric(df["pred_hi_97.5"], errors="coerce").to_numpy(dtype=float)

    good = np.isfinite(x) & np.isfinite(y) & np.isfinite(lo) & np.isfinite(hi)
    x, y, lo, hi = x[good], y[good], lo[good], hi[good]
    idx = np.argsort(x)
    return x[idx], y[idx], lo[idx], hi[idx]


def find_boot_matrix(curve_fp: Path):
    p = Path(str(curve_fp).replace("_curve_boot.csv", "_boot_raw_curve_matrix.csv"))
    if p.exists():
        return p
    cands = list(curve_fp.parent.glob("*boot_raw_curve_matrix.csv"))
    return cands[0] if cands else None


def delta_from_curve(curve_fp: Path, step: float = 5.0):
    x, y, lo, hi = load_curve(curve_fp)
    x0 = float(np.nanmedian(x))
    x1 = min(float(np.nanmax(x)), x0 + step)
    if x1 <= x0:
        x0, x1 = float(np.nanmin(x)), float(np.nanmax(x))
    d = interp(x, y, x1) - interp(x, y, x0)

    mat_fp = find_boot_matrix(curve_fp)
    if mat_fp is not None and mat_fp.exists():
        mat = pd.read_csv(mat_fp).to_numpy(dtype=float)
        if mat.shape[1] != len(x):
            x_mat = np.linspace(float(np.nanmin(x)), float(np.nanmax(x)), mat.shape[1])
        else:
            x_mat = x
        deltas = []
        for row in mat:
            if not np.all(np.isfinite(row)):
                continue
            deltas.append(interp(x_mat, row, x1) - interp(x_mat, row, x0))
        if len(deltas) >= 10:
            arr = np.array(deltas, dtype=float)
            return d, float(np.quantile(arr, 0.025)), float(np.quantile(arr, 0.975)), len(arr)

    se0 = max((interp(x, hi, x0) - interp(x, lo, x0)) / (2 * 1.96), 1e-9)
    se1 = max((interp(x, hi, x1) - interp(x, lo, x1)) / (2 * 1.96), 1e-9)
    se = np.sqrt(se0**2 + se1**2)
    return d, float(d - 1.96 * se), float(d + 1.96 * se), 0


def parse_n_from_name(name: str):
    m = re.search(r"_nSweep_(\d+)_boot\d+_rowreplace$", name)
    return int(m.group(1)) if m else None


def pick_curve(result_dir: Path, ch: str):
    cands = sorted(result_dir.rglob(f"{ch}_ET_CO2_*_slice_median_curve_boot.csv"))
    return cands[0] if cands else None


CHANNEL_MAP = {
    "rSO2_Ch1": "Left SctO₂ (%)",
    "rSO2_Ch2": "Right SctO₂ (%)",
    "rSO2_Ch3": "SftO₂ (%)"
}

CHANNEL_COLORS = {
    "rSO2_Ch1": "#1f77b4",
    "rSO2_Ch2": "#2ca02c",
    "rSO2_Ch3": "#d62728"
}


def plot_delta(df_ok: pd.DataFrame, out_png: Path, out_pdf: Path):
    plt.figure(figsize=(7.5, 4.2))
    ax = plt.gca()
    
    # Apply clean classic styles
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#616161')
    ax.spines['bottom'].set_color('#616161')
    ax.tick_params(colors='#616161', which='both', width=0.8)
    
    for ch in CHANNELS:
        d = df_ok[df_ok["ycol"] == ch].sort_values("sample_size")
        if d.empty:
            continue
        x = d["sample_size"].to_numpy(dtype=float)
        y = d["delta_rso2_plus5"].to_numpy(dtype=float)
        lo = d["delta_ci_lo"].to_numpy(dtype=float)
        hi = d["delta_ci_hi"].to_numpy(dtype=float)
        yerr = np.vstack([y - lo, hi - y])
        
        plt.errorbar(
            x, y, yerr=yerr, 
            marker="o", linestyle="-", capsize=3, linewidth=1.5, 
            color=CHANNEL_COLORS[ch], 
            label=CHANNEL_MAP[ch]
        )

    plt.xscale("log")
    plt.axvline(10000, color="#7f7f7f", linestyle="--", linewidth=1.2, alpha=0.85, label="Main Study Baseline (N=10,000)")
    
    plt.xlabel("Subsample Size N (log scale)", fontsize=11, color="black")
    plt.ylabel("Delta rSO₂ (%) for +5 mmHg ET-CO₂", fontsize=11, color="black")
    
    plt.legend(frameon=False, loc="upper right", fontsize=10)
    plt.tight_layout()
    plt.savefig(out_png, dpi=220)
    plt.savefig(out_pdf)
    plt.close()


def plot_curve_overlay(df_ok: pd.DataFrame, out_png: Path, out_pdf: Path):
    picks = [1000, 10000, 100000, 1000000]
    fig, axes = plt.subplots(1, 3, figsize=(13.33, 4.5), sharey=True)

    for i, ch in enumerate(CHANNELS):
        ax = axes[i]
        
        # Apply clean classic styles to each panel
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#616161')
        ax.spines['bottom'].set_color('#616161')
        ax.tick_params(colors='#616161', width=0.8)
        
        dch = df_ok[df_ok["ycol"] == ch].copy()
        if dch.empty:
            ax.set_title(CHANNEL_MAP[ch], fontsize=11, fontweight="bold")
            continue

        ns_avail = sorted(dch["sample_size"].unique())
        draw_ns = [n for n in picks if n in ns_avail]
        if not draw_ns:
            draw_ns = ns_avail[: min(4, len(ns_avail))]

        cmap = plt.get_cmap("viridis")
        for j, n in enumerate(draw_ns):
            one = dch[dch["sample_size"] == n].iloc[0]
            curve_fp = Path(one["curve_fp"])
            x, y, lo, hi = load_curve(curve_fp)
            color = cmap(j / max(1, len(draw_ns) - 1))
            ax.plot(x, y, color=color, linewidth=2.0, label=f"N={n:,}")
            ax.fill_between(x, lo, hi, color=color, alpha=0.10)

        ax.set_title(CHANNEL_MAP[ch], fontsize=12, fontweight="bold", pad=10)
        ax.set_xlabel("ET-CO₂ (mmHg)", fontsize=11, color="black")
        ax.grid(alpha=0.15, linestyle=":")
        if i == 0:
            ax.set_ylabel("Predicted Oxygenation (%)", fontsize=11, color="black")
            
        # Place compact vertical legend on the upper left of EACH panel
        ax.legend(frameon=False, loc="upper left", fontsize=9.0, labelspacing=0.18, handlelength=1.2)

    plt.tight_layout()
    plt.savefig(out_png, dpi=220, bbox_inches="tight")
    plt.savefig(out_pdf, bbox_inches="tight")
    plt.close()


def build_stability_summary(df_ok: pd.DataFrame):
    rows = []
    for ch in CHANNELS:
        dch = df_ok[df_ok["ycol"] == ch].sort_values("sample_size")
        if dch.empty:
            continue
        ref = dch.iloc[-1]
        d10k = dch[dch["sample_size"] == 10000]
        if d10k.empty:
            continue
        d10k = d10k.iloc[0]
        abs_diff = abs(float(d10k["delta_rso2_plus5"]) - float(ref["delta_rso2_plus5"]))
        rel_diff = abs_diff / (abs(float(ref["delta_rso2_plus5"])) + 1e-9)

        # heuristic plateau flag
        plateau = (abs_diff <= 0.2) or (rel_diff <= 0.1)
        rows.append(
            {
                "ycol": ch,
                "n_ref": int(ref["sample_size"]),
                "delta_ref": float(ref["delta_rso2_plus5"]),
                "delta_n10000": float(d10k["delta_rso2_plus5"]),
                "abs_diff_n10000_vs_ref": float(abs_diff),
                "rel_diff_n10000_vs_ref": float(rel_diff),
                "n10000_in_stable_plateau": bool(plateau),
            }
        )
    return pd.DataFrame(rows)


def build_ppt(df_ok: pd.DataFrame, fig1_png: Path, fig2_png: Path, stab: pd.DataFrame, out_pptx: Path):
    if not HAS_PPTX:
        print("[pptx] python-pptx is not installed, skipping PPTX slide generation.")
        return
    from datetime import datetime
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s0, "Model B Sample Size Sensitivity Sweep Summary", size=24)
    add_body(
        s0, 
        f"Sensitivity analysis on subsample size N\n"
        f"Bootstrap resamples: b=50\n"
        f"Date: {datetime.now().strftime('%Y-%m-%d')}",
        top=1.5,
        size=14
    )

    # Slide 1: Delta effect size over different sample sizes
    if fig1_png.exists():
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s1, "ET-CO₂ Effect Size Stability vs Sample Size N", size=20)
        s1.shapes.add_picture(str(fig1_png), Inches(2.91), Inches(1.8), width=Inches(7.5))
        
    # Slide 2: Curve overlays
    if fig2_png.exists():
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s2, "ET-CO₂ Response Curves across Subsample Sizes", size=20)
        s2.shapes.add_picture(str(fig2_png), Inches(0.91), Inches(1.8), width=Inches(11.5))

    # Slide 3: Stability table
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "Stability Comparison: N=10,000 vs Maximum Sample Size N_ref", size=20)
    
    if stab is not None and not stab.empty:
        rows_cnt = len(stab) + 1
        cols_cnt = 7
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(10.33)
        height = Inches(0.4 * rows_cnt)
        
        table_shape = s3.shapes.add_table(rows_cnt, cols_cnt, left, top, width, height)
        table = table_shape.table
        
        headers = ["Channel", "N_ref", "Delta_ref", "Delta_10k", "Abs Diff", "Rel Diff", "Stable?"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            # Format header font
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)
                    run.font.name = "Arial"
                    
        for row_idx, (_, row) in enumerate(stab.iterrows(), start=1):
            ch = row["ycol"]
            n_ref = int(row["n_ref"])
            d_ref = float(row["delta_ref"])
            d_10k = float(row["delta_n10000"])
            abs_diff = float(row["abs_diff_n10000_vs_ref"])
            rel_diff = float(row["rel_diff_n10000_vs_ref"])
            stable = "Yes" if bool(row["n10000_in_stable_plateau"]) else "No"
            
            # Map channel name to publication name
            ch_pub = CHANNEL_MAP.get(ch, ch)
            
            vals = [
                ch_pub,
                f"{n_ref:,}",
                f"{d_ref:.4f}",
                f"{d_10k:.4f}",
                f"{abs_diff:.4f}",
                f"{rel_diff:.2%}",
                stable
            ]
            
            for col_idx, val in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
                        run.font.name = "Arial"
    else:
        add_body(s3, "No stability comparison data available.", top=2.0, size=14)
        
    prs.save(str(out_pptx))
    print(f"Generated summary slides PPTX: {out_pptx}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--result-root", default="/N/project/waveform_mortality/ZhaoZhang/contour_zhao_all_9_15_2025/result")
    ap.add_argument("--out-dir", default="/N/project/waveform_mortality/ZhaoZhang/co2_rso2_repo_20260516/results/modelb_n_sweep_eval")
    ap.add_argument("--step", type=float, default=5.0)
    args = ap.parse_args()

    result_root = Path(args.result_root)
    out_dir = Path(args.out_dir)
    out_tables = out_dir / "tables"
    out_figs = out_dir / "figures"
    out_tables.mkdir(parents=True, exist_ok=True)
    out_figs.mkdir(parents=True, exist_ok=True)

    dirs = sorted(result_root.glob("*overall_modelB_sec1_nSweep_*_boot50_rowreplace"))
    rows = []

    for d in dirs:
        n = parse_n_from_name(d.name)
        if n is None:
            continue
        run_summary_ok = (d / "run_summary.csv").exists()

        for ch in CHANNELS:
            curve = pick_curve(d, ch)
            if curve is None:
                rows.append(
                    {
                        "sample_size": n,
                        "ycol": ch,
                        "status": "missing_curve",
                        "result_dir": str(d),
                        "run_summary_ok": int(run_summary_ok),
                    }
                )
                continue
            delta, lo, hi, n_boot = delta_from_curve(curve, step=args.step)
            rows.append(
                {
                    "sample_size": n,
                    "ycol": ch,
                    "status": "ok",
                    "result_dir": str(d),
                    "run_summary_ok": int(run_summary_ok),
                    "curve_fp": str(curve),
                    "delta_rso2_plus5": delta,
                    "delta_ci_lo": lo,
                    "delta_ci_hi": hi,
                    "n_boot_for_ci": int(n_boot),
                }
            )

    df = pd.DataFrame(rows).sort_values(["sample_size", "ycol"]) if rows else pd.DataFrame()
    sum_fp = out_tables / "modelB_n_sweep_etco2_delta_summary.csv"
    df.to_csv(sum_fp, index=False)

    if df.empty:
        print(f"no n-sweep result dirs found under: {result_root}")
        print(f"summary: {sum_fp}")
        return

    df_ok = df[df["status"] == "ok"].copy()
    if not df_ok.empty:
        fig1_png = out_figs / "modelB_n_sweep_delta_plus5_by_channel.png"
        fig1_pdf = out_figs / "modelB_n_sweep_delta_plus5_by_channel.pdf"
        plot_delta(df_ok, fig1_png, fig1_pdf)
        
        fig2_png = out_figs / "modelB_n_sweep_curve_overlay_by_channel.png"
        fig2_pdf = out_figs / "modelB_n_sweep_curve_overlay_by_channel.pdf"
        plot_curve_overlay(df_ok, fig2_png, fig2_pdf)

        stab = build_stability_summary(df_ok)
        stab_fp = out_tables / "modelB_n_sweep_stability_summary.csv"
        stab.to_csv(stab_fp, index=False)
        print(f"stability: {stab_fp}")
        
        # Build PPTX slides
        out_pptx = out_figs / "modelB_n_sweep_etco2_stability_summary.pptx"
        build_ppt(df_ok, fig1_png, fig2_png, stab, out_pptx)

    print(f"summary: {sum_fp}")
    print(f"figures: {out_figs}")


if __name__ == "__main__":
    main()

