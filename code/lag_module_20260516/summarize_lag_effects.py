#!/usr/bin/env python3
"""
v1.1 - 滞后效应结果汇总与 PPTX 报告生成脚本
改动：
  - 支持模型参数动态获取（Model A / Model B）。
  - 支持输出动态文件名称与对应的子采样和自助抽样参数。
  - 增加使用 python-pptx 自动将结果折线图与曲线对比图生成专业级 PPTX 幻灯片功能。
基于：原始 summarize_lag_effects.py
解决问题：支持动态模型参数配置，并实现汇报 PPT 自动化生成。
"""
import argparse
from pathlib import Path
import os


import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

# Try importing python-pptx for professional report generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

CHANNELS = ["rSO2_Ch1", "rSO2_Ch2", "rSO2_Ch3"]


def interp(x, y, xq):
    return float(np.interp(float(xq), x, y))


def find_result_dir(result_root: Path, outdir_tag: str):
    cands = sorted(result_root.glob(f"*_{outdir_tag}"))
    if not cands:
        return None
    # Prioritize directories that actually completed and have summaries
    with_summary = [d for d in cands if (d / "run_summary.csv").exists()]
    return with_summary[-1] if with_summary else cands[-1]


def load_curve(curve_fp: Path):
    df = pd.read_csv(curve_fp)
    x = pd.to_numeric(df["x"], errors="coerce").to_numpy(dtype=float)
    y = pd.to_numeric(df["pred_mean"], errors="coerce").to_numpy(dtype=float)
    lo = pd.to_numeric(df["pred_lo_2.5"], errors="coerce").to_numpy(dtype=float)
    hi = pd.to_numeric(df["pred_hi_97.5"], errors="coerce").to_numpy(dtype=float)

    good = np.isfinite(x) & np.isfinite(y) & np.isfinite(lo) & np.isfinite(hi)
    x, y, lo, hi = x[good], y[good], lo[good], hi[good]
    idx = np.argsort(x)
    return x[idx], y[idx], lo[idx], hi[idx]


def find_boot_matrix(curve_fp: Path):
    cand1 = Path(str(curve_fp).replace("_curve_boot.csv", "_boot_raw_curve_matrix.csv"))
    if cand1.exists():
        return cand1
    cands = list(curve_fp.parent.glob("*boot_raw_curve_matrix.csv"))
    return cands[0] if cands else None


def delta_from_curve(curve_fp: Path, step: float = 5.0):
    x, y, lo, hi = load_curve(curve_fp)
    x0 = float(np.nanmedian(x))
    x1 = min(float(np.nanmax(x)), x0 + float(step))
    if x1 <= x0:
        x0, x1 = float(np.nanmin(x)), float(np.nanmax(x))

    d = interp(x, y, x1) - interp(x, y, x0)

    mat_fp = find_boot_matrix(curve_fp)
    if mat_fp is not None and mat_fp.exists():
        mat = pd.read_csv(mat_fp).to_numpy(dtype=float)
        if mat.shape[1] != len(x):
            x_mat = np.linspace(float(np.nanmin(x)), float(np.nanmax(x)), mat.shape[1])
        else:
            x_mat = x
        deltas = []
        for row in mat:
            if not np.all(np.isfinite(row)):
                continue
            deltas.append(interp(x_mat, row, x1) - interp(x_mat, row, x0))
        if len(deltas) >= 10:
            arr = np.array(deltas, dtype=float)
            return d, float(np.quantile(arr, 0.025)), float(np.quantile(arr, 0.975)), x0, x1, int(len(arr))

    se0 = max((interp(x, hi, x0) - interp(x, lo, x0)) / (2 * 1.96), 1e-9)
    se1 = max((interp(x, hi, x1) - interp(x, lo, x1)) / (2 * 1.96), 1e-9)
    se = np.sqrt(se0**2 + se1**2)
    return d, float(d - 1.96 * se), float(d + 1.96 * se), x0, x1, 0


def pick_strongest_lag_nonzero(df: pd.DataFrame):
    d = df[(df["status"] == "ok") & (df["lag_seconds"] > 0)].copy()
    if d.empty:
        return None
    g = d.groupby("lag_seconds", as_index=False)["delta_rso2_plus5"].mean()
    g["score"] = g["delta_rso2_plus5"].abs()
    g = g.sort_values(["score", "lag_seconds"], ascending=[False, True])
    return int(g.iloc[0]["lag_seconds"])


def extract_sample_counts(flow_fp: Path, ycol: str):
    if not flow_fp.exists():
        return np.nan, np.nan, np.nan, np.nan
    d = pd.read_csv(flow_fp)
    d = d[(d["sec"] == 1) & (d["subgroup"] == "All") & (d["ycol"] == ycol)]
    if d.empty:
        return np.nan, np.nan, np.nan, np.nan

    def get_stage(stage_name: str, col: str):
        s = d[d["stage"] == stage_name][col]
        if len(s) == 0:
            return np.nan
        return float(s.iloc[0])

    n_required_rows = get_stage("after_required_etco2_y_nonmissing", "n_rows")
    n_required_patients = get_stage("after_required_etco2_y_nonmissing", "n_patients")
    n_final_rows = get_stage("final_usable_points_strict_etco2_rso2", "n_rows")
    n_final_patients = get_stage("final_usable_points_strict_etco2_rso2", "n_patients")
    return n_required_rows, n_required_patients, n_final_rows, n_final_patients


def plot_lag_delta(df: pd.DataFrame, out_png: Path, out_pdf: Path, model_label: str, subsample: int, boot: int):
    ok = df[df["status"] == "ok"].copy()
    if ok.empty:
        return

    plt.figure(figsize=(8.6, 5.4))
    colors = {
        "rSO2_Ch1": "#1f77b4",
        "rSO2_Ch2": "#ff7f0e",
        "rSO2_Ch3": "#2ca02c",
    }
    for ch in CHANNELS:
        sub = ok[ok["ycol"] == ch].sort_values("lag_seconds")
        if sub.empty:
            continue
        x = sub["lag_seconds"].to_numpy(dtype=float)
        y = sub["delta_rso2_plus5"].to_numpy(dtype=float)
        lo = sub["delta_ci_lo"].to_numpy(dtype=float)
        hi = sub["delta_ci_hi"].to_numpy(dtype=float)
        yerr = np.vstack([y - lo, hi - y])
        plt.errorbar(
            x,
            y,
            yerr=yerr,
            marker="o",
            linestyle="-",
            linewidth=1.5,
            capsize=3,
            color=colors.get(ch, None),
            label=ch,
        )

    plt.axhline(0.0, color="black", linestyle="--", linewidth=1.0)
    plt.xlabel("ET_CO2 lag (seconds)")
    plt.ylabel("Delta rSO2 (+5 mmHg ET_CO2)")
    plt.title(f"Lag sensitivity of ET_CO2 effect ({model_label}, n={subsample}, boot={boot})")
    plt.xticks([0, 30, 60, 120, 180])
    plt.legend(frameon=False)
    plt.tight_layout()
    plt.savefig(out_png, dpi=220)
    plt.savefig(out_pdf)
    plt.close()


def plot_curve_compare(df: pd.DataFrame, lag0: int, lag_best: int, out_png: Path, out_pdf: Path):
    fig, axes = plt.subplots(1, 3, figsize=(14.5, 4.2), sharey=True)
    if not isinstance(axes, np.ndarray):
        axes = np.array([axes])

    for i, ch in enumerate(CHANNELS):
        ax = axes[i]
        sub = df[(df["status"] == "ok") & (df["ycol"] == ch) & (df["lag_seconds"].isin([lag0, lag_best]))]
        for lag_s, color in [(lag0, "#1f77b4"), (lag_best, "#d62728")]:
            one = sub[sub["lag_seconds"] == lag_s]
            if one.empty:
                continue
            curve_fp = Path(one.iloc[0]["curve_fp"])
            x, y, lo, hi = load_curve(curve_fp)
            ax.plot(x, y, color=color, linewidth=2.0, label=f"lag{lag_s}s")
            ax.fill_between(x, lo, hi, color=color, alpha=0.15)

        ax.set_title(ch)
        ax.set_xlabel("ET_CO2 (mmHg)")
        if i == 0:
            ax.set_ylabel("Predicted rSO2")
        ax.grid(alpha=0.2, linewidth=0.6)

    handles, labels = axes[0].get_legend_handles_labels()
    if handles:
        fig.legend(handles, labels, loc="lower center", ncol=2, frameon=False)
    fig.suptitle(f"ET_CO2 slice curves: lag0 vs lag{lag_best}s", y=1.02)
    plt.tight_layout()
    plt.savefig(out_png, dpi=220, bbox_inches="tight")
    plt.savefig(out_pdf, bbox_inches="tight")
    plt.close()


# pptx helper functions
def add_title(slide, text, left=0.5, top=0.2, width=12.0, height=0.6, size=28):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True


def add_body(slide, text, left=0.6, top=1.0, width=12.0, height=1.5, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)


def build_ppt(model_label: str, subsample: int, boot: int, fig1_png: Path, fig2_png: Path, df_ok: pd.DataFrame, out_pptx: Path):
    if not HAS_PPTX:
        print("[pptx] python-pptx is not installed, skipping PPTX slide generation.")
        return

    prs = Presentation()
    
    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s0, "ET_CO2 Lag Sensitivity Analysis Summary", size=24)
    add_body(
        s0, 
        f"Model: {model_label}\n"
        f"Subsample size: n={subsample}\n"
        f"Bootstrap resamples: b={boot}\n"
        f"Date: {datetime.now().strftime('%Y-%m-%d')}",
        top=1.5,
        size=14
    )

    # Slide 1: Delta effect size over different lag settings
    if fig1_png.exists():
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s1, "Lag sensitivity of ET_CO2 effect size (+5 mmHg)", size=20)
        s1.shapes.add_picture(str(fig1_png), Inches(0.5), Inches(0.9), width=Inches(9.0))
        
    # Slide 2: Curve comparisons (lag0 vs lag_best)
    if fig2_png.exists():
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s2, "ET_CO2 response curves: lag0 vs best nonzero lag", size=20)
        s2.shapes.add_picture(str(fig2_png), Inches(0.5), Inches(0.9), width=Inches(9.0))

    # Slide 3: Key numeric values
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "Appendix: Key numeric summaries", size=20)
    
    # Build text summary of numeric outputs
    txt_summary = ""
    for ch in CHANNELS:
        ch_df = df_ok[df_ok["ycol"] == ch].sort_values("lag_seconds")
        txt_summary += f"--- {ch} ---\n"
        for _, row in ch_df.iterrows():
            lag = row["lag_seconds"]
            d = row["delta_rso2_plus5"]
            lo = row["delta_ci_lo"]
            hi = row["delta_ci_hi"]
            patients = int(row["n_final_patients"]) if pd.notna(row["n_final_patients"]) else 0
            txt_summary += f"  Lag {lag:3d}s: Delta={d:5.2f} (95% CI: {lo:5.2f} to {hi:5.2f}) | Patients={patients}\n"
        txt_summary += "\n"
    
    add_body(s3, txt_summary, top=1.0, size=11, height=5.5)
    
    prs.save(str(out_pptx))
    print(f"Generated summary slides PPTX: {out_pptx}")


from datetime import datetime


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument(
        "--workspace",
        default="/N/project/waveform_mortality/ZhaoZhang/co2_rso2_repo_20260516/code/lag_module_20260516",
    )
    ap.add_argument("--matrix", default="lag_run_matrix.csv")
    ap.add_argument(
        "--result-root",
        default="/N/project/waveform_mortality/ZhaoZhang/co2_rso2_repo_20260516/code/analysis_bundle/result",
    )
    ap.add_argument("--step", type=float, default=5.0)
    ap.add_argument("--model", choices=["A", "B"], default="B", help="Model type: A (map_ci_te) or B (map_sv_smooth) (default: B)")
    ap.add_argument("--boot", type=int, default=200, help="Number of bootstrap resamples (default: 200)")
    ap.add_argument("--subsample", type=int, default=10000, help="Subsample size (default: 10000)")
    args = ap.parse_args()

    ws = Path(args.workspace)
    matrix_fp = ws / args.matrix
    result_root = Path(args.result_root)
    out_tables = ws / "output" / "tables"
    out_figs = ws / "output" / "figures"
    out_tables.mkdir(parents=True, exist_ok=True)
    out_figs.mkdir(parents=True, exist_ok=True)

    model_label = "modelA" if args.model == "A" else "modelB"
    matrix = pd.read_csv(matrix_fp)
    rows = []

    for _, r in matrix.iterrows():
        if int(r.get("enabled", 0)) != 1:
            continue
        run_key = str(r["run_key"])
        lag_seconds = int(r["lag_seconds"])
        
        # Dynamically compute tag to find results
        outdir_tag = f"lag{lag_seconds}_{model_label}_n{args.subsample}_b{args.boot}"
        result_dir = find_result_dir(result_root, outdir_tag)

        if result_dir is None:
            for ch in CHANNELS:
                rows.append({
                    "run_key": run_key,
                    "lag_seconds": lag_seconds,
                    "outdir_tag": outdir_tag,
                    "ycol": ch,
                    "status": "missing_result_dir",
                })
            continue

        flow_fp = result_dir / "filter_flow_counts.csv"
        for ch in CHANNELS:
            curve = next(result_dir.rglob(f"{ch}_ET_CO2_*_slice_median_curve_boot.csv"), None)
            if curve is None:
                rows.append({
                    "run_key": run_key,
                    "lag_seconds": lag_seconds,
                    "outdir_tag": outdir_tag,
                    "ycol": ch,
                    "status": "missing_curve",
                    "result_dir": str(result_dir),
                })
                continue

            d, lo, hi, x0, x1, n_boot_ok = delta_from_curve(curve, step=args.step)
            n_required_rows, n_required_patients, n_final_rows, n_final_patients = extract_sample_counts(flow_fp, ch)

            rows.append({
                "run_key": run_key,
                "lag_seconds": lag_seconds,
                "outdir_tag": outdir_tag,
                "ycol": ch,
                "status": "ok",
                "result_dir": str(result_dir),
                "curve_fp": str(curve),
                "x_from_median": x0,
                "x_to_plus5": x1,
                "delta_rso2_plus5": d,
                "delta_ci_lo": lo,
                "delta_ci_hi": hi,
                "n_boot_ok": n_boot_ok,
                "n_required_rows": n_required_rows,
                "n_required_patients": n_required_patients,
                "n_final_rows": n_final_rows,
                "n_final_patients": n_final_patients,
            })

    df = pd.DataFrame(rows)
    bych_fp = out_tables / f"lag_effect_summary_by_channel_{model_label}_n{args.subsample}_b{args.boot}.csv"
    df.to_csv(bych_fp, index=False)

    ok = df[df["status"] == "ok"].copy()
    if ok.empty:
        print(f"No completed lag results found under: {result_root} for model={model_label}, boot={args.boot}")
        print(f"Wrote missing summary list to: {bych_fp}")
        return

    lag0_df = ok[ok["lag_seconds"] == 0][["ycol", "n_final_rows", "n_final_patients"]].rename(
        columns={"n_final_rows": "lag0_n_final_rows", "n_final_patients": "lag0_n_final_patients"}
    )
    ok = ok.merge(lag0_df, on="ycol", how="left")
    ok["rows_loss_vs_lag0"] = ok["lag0_n_final_rows"] - ok["n_final_rows"]

    bych_withloss_fp = out_tables / f"lag_effect_summary_by_channel_with_loss_{model_label}_n{args.subsample}_b{args.boot}.csv"
    ok.to_csv(bych_withloss_fp, index=False)

    lag_tbl = (
        ok.groupby("lag_seconds", as_index=False)
        .agg(
            mean_delta=("delta_rso2_plus5", "mean"),
            min_ci_lo=("delta_ci_lo", "min"),
            max_ci_hi=("delta_ci_hi", "max"),
            mean_final_rows=("n_final_rows", "mean"),
            mean_rows_loss_vs_lag0=("rows_loss_vs_lag0", "mean"),
        )
        .sort_values("lag_seconds")
    )
    lag_tbl_fp = out_tables / f"lag_effect_summary_overall_{model_label}_n{args.subsample}_b{args.boot}.csv"
    lag_tbl.to_csv(lag_tbl_fp, index=False)

    fig1_png = out_figs / f"lag_delta_plus5_by_channel_{model_label}_n{args.subsample}_b{args.boot}.png"
    fig1_pdf = out_figs / f"lag_delta_plus5_by_channel_{model_label}_n{args.subsample}_b{args.boot}.pdf"
    plot_lag_delta(ok, fig1_png, fig1_pdf, model_label, args.subsample, args.boot)

    lag_best = pick_strongest_lag_nonzero(ok)
    if lag_best is None:
        lag_best = 60

    fig2_png = out_figs / f"lag_curve_compare_lag0_vs_lag{lag_best}_{model_label}_n{args.subsample}_b{args.boot}.png"
    fig2_pdf = out_figs / f"lag_curve_compare_lag0_vs_lag{lag_best}_{model_label}_n{args.subsample}_b{args.boot}.pdf"
    plot_curve_compare(ok, lag0=0, lag_best=lag_best, out_png=fig2_png, out_pdf=fig2_pdf)

    # Build PPTX slides
    out_pptx = out_figs / f"lag_effects_summary_{model_label}_n{args.subsample}_b{args.boot}.pptx"
    build_ppt(model_label, args.subsample, args.boot, fig1_png, fig2_png, ok, out_pptx)

    print(f"summary_by_channel: {bych_fp}")
    print(f"summary_by_channel_with_loss: {bych_withloss_fp}")
    print(f"summary_overall: {lag_tbl_fp}")
    print(f"figure1: {fig1_png}")
    print(f"figure2: {fig2_png}")


if __name__ == "__main__":
    main()
